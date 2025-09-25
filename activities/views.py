from rest_framework import generics, permissions
from .models import Lesson
from .serializers import LessonSerializer
from rest_framework.response import Response
from rest_framework import status
from .utils import chunk_text
from .ai_service import call_ai
import os
import PyPDF2
from pptx import Presentation
from rest_framework.parsers import MultiPartParser, FormParser
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import IsAuthenticated
import json
import re

class LessonCreateView(generics.CreateAPIView):
    queryset = Lesson.objects.all()
    serializer_class = LessonSerializer
    permission_classes = [permissions.IsAuthenticated]
    parser_classes = [MultiPartParser, FormParser]  # allow file uploads

    def extract_text_from_pdf(self, file_obj):
        reader = PyPDF2.PdfReader(file_obj)
        text = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        return " ".join(text)

    def extract_text_from_pptx(self, file_obj):
        prs = Presentation(file_obj)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return " ".join(text)

    def _split_for_model(self, content: str, max_chars: int):
        """
        Split content into chunks <= max_chars trying to respect paragraph boundaries.
        """
        if len(content) <= max_chars:
            return [content]
        paragraphs = re.split(r'\n{2,}|\r\n{2,}', content)
        chunks = []
        current = []
        current_len = 0
        for p in paragraphs:
            p = p.strip()
            if not p:
                continue
            # +2 for the blank line joiner
            add_len = len(p) + 2
            if current_len + add_len > max_chars and current:
                chunks.append("\n\n".join(current))
                current = [p]
                current_len = len(p)
            else:
                current.append(p)
                current_len += add_len
        if current:
            chunks.append("\n\n".join(current))
        # Final safety split (hard cut) if any residual > max_chars
        fixed = []
        for c in chunks:
            if len(c) <= max_chars:
                fixed.append(c)
            else:
                for i in range(0, len(c), max_chars):
                    fixed.append(c[i:i+max_chars])
        return fixed

    def create(self, request, *args, **kwargs):
        title = request.data.get("title")
        raw_topic = request.data.get("topic")
        uploaded_file = request.FILES.get("file")

        if not title:
            return Response({"error": "Title required"}, status=status.HTTP_400_BAD_REQUEST)

        # Acquire raw content (text or file)
        if raw_topic:
            content = raw_topic
        elif uploaded_file:
            ext = os.path.splitext(uploaded_file.name)[1].lower()
            if ext == ".pdf":
                content = self.extract_text_from_pdf(uploaded_file)
            elif ext in [".pptx", ".ppt"]:
                content = self.extract_text_from_pptx(uploaded_file)
            else:
                return Response({"error": "Unsupported file type"}, status=status.HTTP_400_BAD_REQUEST)
        else:
            return Response({"error": "Provide either text or file"}, status=status.HTTP_400_BAD_REQUEST)

        content = content.strip()
        if not content:
            return Response({"error": "Empty content"}, status=status.HTTP_400_BAD_REQUEST)

        # Model-safe per-pass character limit (fits prompt + output in model context)
        PER_PASS_LIMIT = 12000  # adjust if needed based on model context
        parts = self._split_for_model(content, PER_PASS_LIMIT)
        total_parts = len(parts)

        simplified_parts = []

        if total_parts == 1:
            # Single pass (original flow)
            explain_prompt = f"""
You are an expert educator. Rewrite the following lesson in clear, simple English suitable for students.

Requirements:
- Preserve logical structure (keep/normalize headings).
- Break long paragraphs into smaller ones (2–4 sentences each).
- Keep key scientific / technical terms but add a short parenthetical explanation the first time they appear.
- Neutral, encouraging tone.
- Keep essential definitions and relationships.
- Do NOT omit important concepts.
- Return ONLY the rewritten lesson text (no extra commentary, no concluding summary section).

Original Lesson:
{parts[0]}
"""
            explained_full = call_ai(explain_prompt).strip()
        else:
            # Multi-pass streaming
            for idx, part in enumerate(parts, start=1):
                if idx == 1:
                    role_instructions = f"""This is Part {idx} of {total_parts}. More parts will follow."""
                    intro_rule = "Do NOT write an overall introduction or conclusion."
                elif idx == total_parts:
                    role_instructions = f"""This is Part {idx} of {total_parts} (final)."""
                    intro_rule = "Do NOT add a final conclusion or summary; just rewrite this part."
                else:
                    role_instructions = f"""This is Part {idx} of {total_parts} (middle)."""
                    intro_rule = "Do NOT add an introduction or conclusion; continue seamlessly."

                part_prompt = f"""
You are an expert educator rewriting a large lesson in sequential parts.

{role_instructions}

Guidelines:
- Preserve existing logical/heading structure (adjust numbering if needed).
- Break long paragraphs (2–4 sentences).
- First occurrence of key technical terms: add short parenthetical explanation.
- Keep definitions, relationships, and important concepts.
- Maintain neutral, encouraging tone.
- {intro_rule}
- Do NOT reference other parts explicitly.
- Return ONLY rewritten content of this part (no extra commentary).

Original Part Text:
{part}
"""
                rewritten = call_ai(part_prompt).strip()
                # Light cleanup to reduce duplicate leading headings like "Introduction"
                simplified_parts.append(rewritten)

            # Merge parts then run a normalization pass to remove duplicated headings/intro snippets
            merged = "\n\n".join(simplified_parts)
            normalize_prompt = f"""
You are an expert editor.

Task:
Merge the following rewritten lesson fragments into a single cohesive lesson.

Rules:
- Remove duplicate introductions or repeated headings.
- Keep a single coherent flow of sections.
- Preserve all concepts.
- Keep parenthetical explanations already present.
- Do NOT add a concluding summary section.
- Return ONLY the cleaned unified lesson.

Fragments:
{merged}
"""
            explained_full = call_ai(normalize_prompt).strip()

        # Generate quiz from unified simplified lesson
        quiz_prompt = f"""
Create a comprehensive multiple-choice quiz that tests deep understanding of the lesson below.

Requirements:
- 12–20 questions.
- Each question has 4 options labeled A, B, C, D.
- After each question (or its options), clearly mark the correct option using: Correct: B
- Mix question types: definition, application, scenario, cause-effect, comparison.
- Avoid trivial recall; emphasize reasoning.
- Do not repeat wording verbatim.
- Keep formatting clean and consistent.

Lesson:
{explained_full}
"""
        quiz = call_ai(quiz_prompt).strip()

        # Chunk final simplified lesson for slides
        slides = chunk_text(explained_full, max_length=400) or [explained_full]

        lesson = Lesson.objects.create(
            title=title,
            topic=slides,
            quiz=quiz,
            created_by=request.user,
        )
        serializer = LessonSerializer(lesson)
        return Response(serializer.data, status=status.HTTP_201_CREATED)





# Retrieve + Update + Delete a lesson
class LessonDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = Lesson.objects.all()
    serializer_class = LessonSerializer
    permission_classes = [permissions.IsAuthenticated]


# List all lessons
class LessonsListView(generics.ListAPIView):
    queryset = Lesson.objects.all()
    serializer_class = LessonSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        serializer = self.get_serializer(queryset, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def grade_quiz(request, pk):
    try:
        lesson = Lesson.objects.get(pk=pk)
    except Lesson.DoesNotExist:
        return Response({"error": "Lesson not found"}, status=status.HTTP_404_NOT_FOUND)
    
    quiz_data = request.data
    questions = quiz_data.get('questions', [])
    
    if not questions:
        return Response({"error": "No questions provided"}, status=status.HTTP_400_BAD_REQUEST)
    
    # Calculate basic score
    correct_count = 0
    question_results = []
    
    for question_data in questions:
        user_answer = (question_data.get('userAnswer') or '').strip().upper().replace(')', '')
        correct_answer = (question_data.get('correctAnswer') or '').strip().upper().replace(')', '')
        is_correct = user_answer == correct_answer
        
        if is_correct:
            correct_count += 1
            
        question_results.append({
            'question': question_data.get('question', ''),
            'userAnswer': user_answer,
            'correctAnswer': correct_answer,
            'isCorrect': is_correct,
            'explanation': ''  # Will be filled by AI
        })
    
    total_questions = len(questions)
    percentage = (correct_count / total_questions) * 100 if total_questions > 0 else 0
    
    # Generate AI feedback
    quiz_summary = f"""
    Student Quiz Results:
    - Score: {correct_count}/{total_questions} ({percentage:.1f}%)
    - Questions and Answers:
    """
    
    for i, q in enumerate(question_results):
        quiz_summary += f"""
    {i+1}. {q['question']}
       Student answered: {q['userAnswer']}
       Correct answer: {q['correctAnswer']}
       Result: {'Correct' if q['isCorrect'] else 'Incorrect'}
    """
    
    feedback_prompt = f"""
    Based on the quiz results below, provide:
    1. Encouraging feedback (2-3 sentences)
    2. Areas for improvement if score < 70%
    3. Congratulations if score >= 70%

    Keep it positive, educational and short.

    {quiz_summary}
    """
    
    try:
        ai_feedback = call_ai(feedback_prompt)
        
        # Generate explanations for incorrect answers
        for question_result in question_results:
            if not question_result['isCorrect']:
                explanation_prompt = f"""
                Question: {question_result['question']}
                Correct Answer: {question_result['correctAnswer']}
                Student Answer: {question_result['userAnswer']}
                
                Provide a brief explanation (1-2 sentences) of why the correct answer is right.
                """
                explanation = call_ai(explanation_prompt)
                question_result['explanation'] = explanation.strip()
        
    except Exception as e:
        ai_feedback = f"Great job completing the quiz! You scored {correct_count} out of {total_questions} questions correctly."
    
    result = {
        'score': correct_count,
        'totalQuestions': total_questions,
        'percentage': round(percentage, 1),
        'feedback': ai_feedback,
        'questionResults': question_results
    }
    
    return Response(result, status=status.HTTP_200_OK)
